"""内容生成服务 - 摘要/提炼/报告/PPT/格式转换"""
import os, logging, re
from typing import List, Dict, Optional
from app.core.llm_client import llm_client
from app.prompts.templates import PromptManager
from app.services.vector_store import VectorStoreService
from app.utils.file_parser import FileParser

logger = logging.getLogger(__name__)

class ContentGenerationService:
    @staticmethod
    def generate_summary(file_id, content, max_length=300, file_path=None, file_content_base64=None, file_type=None):
        from app.services.file_understanding import FileUnderstandingService
        return FileUnderstandingService.generate_summary(file_id, content, max_length, file_path=file_path, file_content_base64=file_content_base64, file_type=file_type)

    @staticmethod
    def extract_content(file_id, content, file_path=None, file_content_base64=None, file_type=None):
        content = FileParser.ensure_text(content, file_path, file_type, file_content_base64)
        truncated = content[:8000] if len(content) > 8000 else content
        messages = PromptManager.format("extract", content=truncated)
        response = llm_client.chat(messages, module="content_generation", file_id=file_id, temperature=0.5)
        return {"file_id": file_id, "content": response.content, "key_points": [], "format": "markdown"}

    @staticmethod
    def generate_report(file_ids, topic, contents=None, file_paths=None, file_contents_base64=None):
        all_content = []
        # 优先使用传入的文件内容（可能已经过 FileParser 处理）
        if contents:
            for fid in file_ids[:10]:
                c = contents.get(fid, "")
                if c:
                    fp = file_paths.get(fid) if file_paths else None
                    b64 = file_contents_base64.get(fid) if file_contents_base64 else None
                    c = FileParser.ensure_text(c, fp, file_content_base64=b64)
                    all_content.append(f"[文件: {fid}]\n{c[:5000]}")
        # 内容为空时回退到向量检索
        if not all_content:
            for fid in file_ids[:10]:
                results = VectorStoreService.search(topic, file_ids=[fid], top_k=3)
                for r in results: all_content.append(f"[文件: {fid}]\n{r.get('chunk_text', '')}")
        combined = "\n\n---\n\n".join(all_content)
        if not combined.strip(): combined = f"关于「{topic}」的相关素材分析..."
        if len(combined) > 8000:
            combined = combined[:8000]
            lb = max(combined.rfind("\n\n"), combined.rfind("。"))
            if lb > 4000: combined = combined[:lb+1]
        messages = PromptManager.format("report", combined_content=combined, topic=topic)
        response = llm_client.chat(messages, module="content_generation", temperature=0.7, max_tokens=4096)
        return {"file_id": "", "content": response.content, "key_points": [], "format": "markdown"}

    @staticmethod
    def generate_ppt(file_ids, topic, contents=None, file_paths=None, file_contents_base64=None):
        all_content = []
        # 优先使用传入的文件内容
        if contents:
            for fid in file_ids[:10]:
                c = contents.get(fid, "")
                if c:
                    fp = file_paths.get(fid) if file_paths else None
                    b64 = file_contents_base64.get(fid) if file_contents_base64 else None
                    c = FileParser.ensure_text(c, fp, file_content_base64=b64)
                    all_content.append(c[:4000])
        # 内容为空时回退到向量检索
        if not all_content:
            for fid in file_ids[:10]:
                results = VectorStoreService.search(topic, file_ids=[fid], top_k=2)
                for r in results: all_content.append(r.get("chunk_text", ""))
        combined = "\n\n".join(all_content)
        if not combined.strip(): combined = f"关于「{topic}」的演示素材..."
        if len(combined) > 6000: combined = combined[:6000]
        messages = PromptManager.format("ppt", combined_content=combined, topic=topic)
        response = llm_client.chat(messages, module="content_generation", temperature=0.7, max_tokens=4096)
        ppt_file_id = ""
        try:
            from pptx import Presentation; from pptx.util import Inches, Pt; from pptx.enum.text import PP_ALIGN
            prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
            p = txBox.text_frame.paragraphs[0]; p.text = topic; p.font.size = Pt(36); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
            for section in response.content.split("## 第")[1:]:
                lines = section.strip().split("\n")
                if not lines: continue
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = lines[0].strip()
                for sep in ["：", ":"]:
                    if sep in title: title = title.split(sep, 1)[-1]
                slide.shapes.title.text = title.strip()
                body = slide.shapes.placeholders[1]; tf = body.text_frame; tf.clear()
                for line in lines[1:]:
                    line = line.strip()
                    if line.startswith("-"):
                        p = tf.add_paragraph(); p.text = line.lstrip("- ").strip(); p.font.size = Pt(16)
                    elif line.startswith("> 备注："):
                        slide.notes_slide.notes_text_frame.text = line.replace("> 备注：", "").strip()
            output_dir = os.path.join(os.path.dirname(__file__), "..", "..", "logs", "generated")
            os.makedirs(output_dir, exist_ok=True)
            ppt_file_id = os.path.join(output_dir, f"ppt_{topic[:20]}.pptx")
            prs.save(ppt_file_id); logger.info(f"PPTX saved: {ppt_file_id}")
        except Exception as e: logger.warning(f"PPTX creation failed: {e}")
        return {"file_id": ppt_file_id, "content": response.content, "key_points": [], "format": "pptx" if ppt_file_id else "markdown"}

    @staticmethod
    def convert_format(file_id, content, target_format, source_format=None, file_path=None, file_content_base64=None, file_type=None):
        """Convert content between document formats.
        
        Supports: txt -> docx, md -> docx, txt -> pdf (as markdown), md -> pdf (as markdown)
        """
        content = FileParser.ensure_text(content, file_path, file_type, file_content_base64)
        if not content or not content.strip():
            return {"file_id": file_id, "content": "", "key_points": [], "format": target_format, "error": "Empty content"}

        # Auto-detect source format if not provided
        if not source_format:
            source_format = "txt"

        result_content = content
        result_format = target_format

        try:
            if source_format == "md" and target_format == "docx":
                result_content = ContentGenerationService._md_to_docx(content)
            elif source_format == "txt" and target_format == "docx":
                result_content = ContentGenerationService._txt_to_docx(content, file_id)
            elif target_format == "pdf":
                # 以markdown格式输出，标记为pdf格式（实际渲染需要前端或服务端PDF库）
                result_content = content
                result_format = "markdown"
            else:
                # 未支持的转换，直接返回原内容
                result_content = content
                result_format = source_format
        except Exception as e:
            logger.error(f"Format conversion failed: {e}")
            return {"file_id": file_id, "content": content, "key_points": [], "format": source_format, "error": str(e)}

        return {"file_id": file_id, "content": result_content, "key_points": [], "format": result_format}

    @staticmethod
    def _txt_to_docx(text, file_id):
        """Convert plain text to DOCX using python-docx"""
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()
        
        # Set default font
        style = doc.styles['Normal']
        font = style.font
        font.name = 'SimSun'
        font.size = Pt(11)

        # Add title from file name
        title = os.path.splitext(os.path.basename(file_id))[0] if file_id else "Document"
        heading = doc.add_heading(title, level=1)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Process paragraphs (split by double newlines for paragraphs)
        paragraphs = re.split(r'\n\s*\n', text)
        for para_text in paragraphs:
            para_text = para_text.strip()
            if not para_text:
                continue
            # Check if it looks like a heading
            if para_text.isupper() and len(para_text) < 100:
                doc.add_heading(para_text.capitalize(), level=2)
            else:
                doc.add_paragraph(para_text)

        # Save to temp file and read back as base64
        output_dir = os.environ.get('CONVERT_OUTPUT_DIR', '/tmp/converted')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"converted_{os.urandom(4).hex()}.docx")
        doc.save(output_path)
        
        with open(output_path, 'rb') as f:
            import base64
            result = base64.b64encode(f.read()).decode('utf-8')
        
        os.remove(output_path)
        return result

    @staticmethod
    def _md_to_docx(markdown_text):
        """Convert Markdown to DOCX using python-docx + markdown"""
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()
        style = doc.styles['Normal']
        font = style.font
        font.name = 'SimSun'
        font.size = Pt(11)

        lines = markdown_text.split('\n')
        i = 0
        while i < len(lines):
            line = lines[i].strip()

            if not line:
                i += 1
                continue

            # Headings
            if line.startswith('### '):
                doc.add_heading(line[4:], level=3)
            elif line.startswith('## '):
                doc.add_heading(line[3:], level=2)
            elif line.startswith('# '):
                doc.add_heading(line[2:], level=1)
            # Unordered list
            elif line.startswith('- ') or line.startswith('* '):
                p = doc.add_paragraph(line[2:], style='List Bullet')
            # Ordered list
            elif re.match(r'^\d+\.\s', line):
                content = re.sub(r'^\d+\.\s', '', line)
                p = doc.add_paragraph(content, style='List Number')
            # Blockquote
            elif line.startswith('> '):
                p = doc.add_paragraph(line[2:])
                p.paragraph_format.left_indent = Inches(0.5)
            # Code block
            elif line.startswith('```'):
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith('```'):
                    code_lines.append(lines[i])
                    i += 1
                code_text = '\n'.join(code_lines)
                p = doc.add_paragraph(code_text)
                run = p.runs[0] if p.runs else p.add_run(code_text)
                run.font.name = 'Courier New'
                run.font.size = Pt(9)
            # Horizontal rule
            elif line in ('---', '***', '___'):
                doc.add_paragraph('_' * 40)
            # Regular paragraph
            else:
                # Check for inline formatting
                text = line
                p = doc.add_paragraph(text)

            i += 1

        # Save and return base64
        output_dir = os.environ.get('CONVERT_OUTPUT_DIR', '/tmp/converted')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"converted_md_{os.urandom(4).hex()}.docx")
        doc.save(output_path)

        with open(output_path, 'rb') as f:
            import base64
            result = base64.b64encode(f.read()).decode('utf-8')

        os.remove(output_path)
        return result

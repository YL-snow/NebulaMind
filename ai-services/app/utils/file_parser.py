"""文件解析器 - PDF/Word/Excel/PPT/图片/TXT"""
import io
import logging

logger = logging.getLogger(__name__)


class FileParser:
    # 二进制文件无法提取有效文本时的固定标记，后端据此触发视觉模型兜底
    NO_EXTRACTABLE_TEXT = "[NO_EXTRACTABLE_TEXT]"

    # 需要解析的二进制文件扩展名
    BINARY_EXTENSIONS = {"pdf", "docx", "doc", "xlsx", "xls", "pptx", "ppt",
                         "jpg", "jpeg", "png", "gif", "bmp", "tiff", "tif", "webp"}

    @staticmethod
    def _ext_of(file_path, file_type):
        if file_type:
            return file_type.lower()
        if file_path and "." in file_path:
            return file_path.rsplit(".", 1)[-1].lower()
        return ""

    @staticmethod
    def _is_binary(ext):
        return (ext or "").lower() in FileParser.BINARY_EXTENSIONS

    @staticmethod
    def _binary_no_text_message():
        return FileParser.NO_EXTRACTABLE_TEXT + " 该文件未提取到可分析的文本，可能为图片型文档或空白文件。"

    @staticmethod
    def ensure_text(content, file_path=None, file_type=None, file_content_base64=None):
        """确保获取到真实的文本内容。

        优先级：
        1. file_content_base64 - base64 编码的原始文件字节（兼容所有存储后端）
        2. file_path - 从共享文件系统中解析二进制文件
        3. content - 原始文本（用于 txt/md 等纯文本格式）

        当文件是二进制格式（pdf/docx/xlsx/pptx）时，使用 FileParser 提取纯文本。
        """
        # 优先使用 base64 编码的原始字节
        if file_content_base64:
            try:
                import base64
                raw_bytes = base64.b64decode(file_content_base64)
                parsed = FileParser.parse(
                    file_path or "file",
                    file_content=raw_bytes,
                    file_type=file_type,
                )
                if parsed and len(parsed.strip()) > 2:
                    logger.info(f"FileParser extracted {len(parsed)} chars from base64 content")
                    return parsed
                logger.warning("FileParser returned empty/short content from base64")
                if FileParser._is_binary(FileParser._ext_of(file_path, file_type)):
                    return FileParser._binary_no_text_message()
                return parsed or content
            except Exception as e:
                logger.warning(f"FileParser failed for base64 content: {e}")
                if FileParser._is_binary(FileParser._ext_of(file_path, file_type)):
                    return FileParser._binary_no_text_message()
                return content

        # 其次尝试从文件路径解析
        if file_path:
            ext = FileParser._ext_of(file_path, file_type)
            if ext in FileParser.BINARY_EXTENSIONS:
                try:
                    parsed = FileParser.parse(file_path, file_type=file_type)
                    if parsed and len(parsed.strip()) > 2:
                        logger.info(f"FileParser extracted {len(parsed)} chars from {file_path}")
                        return parsed
                    logger.warning(f"FileParser returned empty/short content for {file_path}")
                    return FileParser._binary_no_text_message()
                except Exception as e:
                    logger.warning(f"FileParser failed for {file_path}: {e}")
                    return FileParser._binary_no_text_message()
        return content

    @staticmethod
    def parse(file_path, file_content=None, mime_type=None, file_type=None):
        ext = FileParser._ext_of(file_path, file_type)
        parsers = {
            "pdf": FileParser._parse_pdf, "docx": FileParser._parse_docx,
            "doc": FileParser._parse_docx, "xlsx": FileParser._parse_xlsx,
            "xls": FileParser._parse_xlsx, "pptx": FileParser._parse_pptx,
            "ppt": FileParser._parse_pptx,
            # 图片格式 - 使用 OCR
            "jpg": FileParser._parse_image, "jpeg": FileParser._parse_image,
            "png": FileParser._parse_image, "gif": FileParser._parse_image,
            "bmp": FileParser._parse_image, "tiff": FileParser._parse_image,
            "tif": FileParser._parse_image, "webp": FileParser._parse_image,
        }
        # 也通过 mime_type 判断
        if mime_type and ext not in parsers:
            mime_map = {
                "image/jpeg": FileParser._parse_image, "image/png": FileParser._parse_image,
                "image/gif": FileParser._parse_image, "image/bmp": FileParser._parse_image,
                "image/tiff": FileParser._parse_image, "image/webp": FileParser._parse_image,
            }
            parser = mime_map.get(mime_type, FileParser._parse_txt)
        else:
            parser = parsers.get(ext, FileParser._parse_txt)
        try:
            if file_content:
                return parser(file_content)
            else:
                with open(file_path, "rb") as f:
                    return parser(f.read())
        except Exception as e:
            logger.error(f"Parse failed {file_path}: {e}")
            return ""

    @staticmethod
    def _parse_pdf(content):
        try:
            from PyPDF2 import PdfReader
            text = "\n\n".join(p.extract_text() or "" for p in PdfReader(io.BytesIO(content)).pages)
            if text and text.strip():
                return text
        except Exception:
            pass
        return FileParser._parse_pdf_with_maas_doc_parser(content)

    @staticmethod
    def _parse_pdf_with_maas_doc_parser(content):
        """本地 PDF 文本提取为空时，使用元景高精度文档解析接口（支持扫描版 OCR）"""
        try:
            import httpx
            from config import settings
            url = settings.maas_doc_parser_url
            api_key = settings.openai_api_key
            if not api_key:
                logger.warning("MaaS doc parser API key missing")
                return ""
            files = {
                "file": ("document.pdf", content, "application/pdf"),
                "file_name": (None, "document.pdf"),
                "extract_image_content": (None, "1"),
            }
            headers = {"Authorization": "Bearer " + api_key}
            with httpx.Client(timeout=90.0, verify=False) as client:
                resp = client.post(url, headers=headers, files=files)
            if resp.status_code != 200:
                logger.warning(f"MaaS doc parser failed: HTTP {resp.status_code} {resp.text[:200]}")
                return ""
            data = resp.json()
            if str(data.get("code")) != "200":
                logger.warning(f"MaaS doc parser returned error: {resp.text[:300]}")
                return ""
            return data.get("content", "")
        except Exception as e:
            logger.warning(f"MaaS doc parser unavailable: {e}")
            return ""

    @staticmethod
    def _parse_docx(content):
        try:
            from docx import Document
            return "\n\n".join(p.text for p in Document(io.BytesIO(content)).paragraphs if p.text.strip())
        except Exception:
            return ""

    @staticmethod
    def _parse_xlsx(content):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), data_only=True)
            texts = []
            for sn in wb.sheetnames:
                texts.append(f"## Sheet: {sn}")
                texts.append("\n".join("\t".join(str(c or "") for c in row)
                                       for row in wb[sn].iter_rows(values_only=True)))
            return "\n\n".join(texts)
        except Exception:
            return ""

    @staticmethod
    def _parse_pptx(content):
        try:
            from pptx import Presentation
            texts = []
            for i, slide in enumerate(Presentation(io.BytesIO(content)).slides, 1):
                st = [f"## Slide {i}"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                st.append(para.text)
                texts.append("\n".join(st))
            return "\n\n".join(texts)
        except Exception:
            return ""

    @staticmethod
    def _parse_image(content):
        """OCR 图片文字识别：Tesseract -> RapidOCR -> 视觉模型兜底标记"""
        text = FileParser._ocr_with_tesseract(content)
        if text:
            return text
        text = FileParser._ocr_with_rapidocr(content)
        if text:
            return text
        return "[图片文件 - OCR未识别到文字，可使用AI视觉模型兜底]"

    @staticmethod
    def _ocr_with_tesseract(content):
        try:
            from PIL import Image
            import pytesseract
            image = Image.open(io.BytesIO(content))
            try:
                text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            except Exception:
                text = pytesseract.image_to_string(image, lang='eng')
            result = text.strip()
            if result:
                logger.info(f"Tesseract OCR extracted {len(result)} chars from image")
            return result
        except Exception as e:
            logger.warning(f"Tesseract OCR unavailable: {e}")
            return ""

    @staticmethod
    def _ocr_with_rapidocr(content):
        """RapidOCR 为纯 Python 依赖，不依赖系统安装的 Tesseract 可执行文件"""
        try:
            from rapidocr_onnxruntime import RapidOCR
            ocr = RapidOCR()
            result = ocr(content)
            if isinstance(result, tuple):
                result = result[0] if result else None
            lines = []
            for item in result or []:
                if isinstance(item, (list, tuple)) and len(item) >= 2 and item[1]:
                    lines.append(str(item[1]))
                elif isinstance(item, str) and item.strip():
                    lines.append(item.strip())
            text = "\n".join(lines).strip()
            if text:
                logger.info(f"RapidOCR extracted {len(text)} chars from image")
            return text
        except Exception as e:
            logger.warning(f"RapidOCR unavailable: {e}")
            return ""

    @staticmethod
    def _parse_txt(content):
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return content.decode("gbk")
            except UnicodeDecodeError:
                return content.decode("utf-8", errors="replace")
